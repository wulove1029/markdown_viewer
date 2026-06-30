"""Tests for Markdown -> PowerPoint (.pptx) export.

Pure logic + a round-trip: generate a deck and read it back with python-pptx.
"""

import pytest

pytest.importorskip("pptx")

from pptx import Presentation

from app.pptx_export import (
    Code,
    Heading,
    Image,
    ListBlock,
    Para,
    Table,
    export_markdown_to_pptx,
    parse_elements,
    split_into_slides,
)


def _slide_text(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
    return "\n".join(out)


# ------------------------------ parsing ------------------------------
def test_parse_basic_elements():
    els = parse_elements(
        "# Title\n\nA paragraph.\n\n- a\n- b\n\n```py\nprint(1)\n```\n"
    )
    kinds = [type(e).__name__ for e in els]
    assert kinds == ["Heading", "Para", "ListBlock", "Code"]
    assert els[0].level == 1
    assert els[3].lang == "py" and "print(1)" in els[3].text


def test_parse_table_and_nested_list():
    els = parse_elements(
        "| A | B |\n| - | - |\n| 1 | 2 |\n\n- top\n    - nested\n"
    )
    table = next(e for e in els if isinstance(e, Table))
    assert table.header == ["A", "B"]
    assert table.rows == [["1", "2"]]
    lst = next(e for e in els if isinstance(e, ListBlock))
    assert [i.level for i in lst.items] == [0, 1]


# --------------------------- slide splitting ---------------------------
def test_split_by_hr_when_present():
    els = parse_elements("# One\n\nintro\n\n---\n\n# Two\n\nbody\n")
    slides = split_into_slides(els)
    assert len(slides) == 2
    assert "".join(r.text for r in slides[0].title) == "One"
    assert "".join(r.text for r in slides[1].title) == "Two"


def test_split_by_h2_when_no_hr():
    md = "# Deck\n\nlead\n\n## Section A\n\na body\n\n## Section B\n\nb body\n"
    slides = split_into_slides(parse_elements(md))
    titles = ["".join(r.text for r in s.title) if s.title else None for s in slides]
    # H1 starts the first (title) slide, each H2 starts its own.
    assert titles == ["Deck", "Section A", "Section B"]


def test_split_single_slide_without_headings_or_hr():
    slides = split_into_slides(parse_elements("just text\n\nmore text\n"))
    assert len(slides) == 1
    assert slides[0].title is None


# ----------------------------- round trip -----------------------------
def test_export_roundtrip_slide_count_and_titles(tmp_path):
    md = "# Intro\n\nlead\n\n## First\n\nhello\n\n## Second\n\nworld\n"
    out = tmp_path / "deck.pptx"
    count = export_markdown_to_pptx(md, out)
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == count == 3
    texts = [_slide_text(s) for s in prs.slides]
    assert "Intro" in texts[0]
    assert "First" in texts[1] and "hello" in texts[1]
    assert "Second" in texts[2] and "world" in texts[2]


def test_export_code_table_list_present(tmp_path):
    md = (
        "## Slide\n\n"
        "- alpha\n- beta\n\n"
        "```c\nint x = 1;\n```\n\n"
        "| H1 | H2 |\n| - | - |\n| a | b |\n"
    )
    out = tmp_path / "rich.pptx"
    export_markdown_to_pptx(md, out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    text = _slide_text(slide)
    assert "alpha" in text and "beta" in text
    assert "int x = 1;" in text
    has_table = any(shape.has_table for shape in slide.shapes)
    assert has_table


def test_export_empty_document_makes_one_slide(tmp_path):
    out = tmp_path / "empty.pptx"
    count = export_markdown_to_pptx("", out)
    assert count == 1
    assert len(Presentation(str(out)).slides) == 1


@pytest.fixture
def png_fixture(tmp_path):
    import pymupdf

    pix = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 40, 24))
    pix.clear_with(210)
    p = tmp_path / "frag.png"
    pix.save(str(p))
    return str(p)


def _has_picture(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def test_image_provider_embeds_picture_for_mermaid(tmp_path, png_fixture):
    md = "## Diagram\n\n```mermaid\ngraph TD; A-->B;\n```\n"
    out = tmp_path / "m.pptx"
    calls = []

    def provider(kind, source):
        calls.append((kind, source))
        return png_fixture if kind == "mermaid" else None

    export_markdown_to_pptx(md, out, image_provider=provider)
    assert calls and calls[0][0] == "mermaid" and "A-->B" in calls[0][1]
    slide = Presentation(str(out)).slides[0]
    assert _has_picture(slide)
    # The mermaid source must NOT also appear as a code box.
    assert "graph TD" not in _slide_text(slide)


def test_image_provider_none_falls_back_to_source_box(tmp_path):
    md = "## Math\n\n$$\n\\int_0^1 x\\,dx\n$$\n"
    out = tmp_path / "math.pptx"
    export_markdown_to_pptx(md, out, image_provider=lambda kind, src: None)
    slide = Presentation(str(out)).slides[0]
    assert not _has_picture(slide)
    assert "int_0^1" in _slide_text(slide)  # labelled source box


def test_no_provider_keeps_source_box(tmp_path):
    md = "## D\n\n```mermaid\ngraph LR; X-->Y;\n```\n"
    out = tmp_path / "noprov.pptx"
    export_markdown_to_pptx(md, out)  # default image_provider=None
    slide = Presentation(str(out)).slides[0]
    assert not _has_picture(slide)
    assert "graph LR" in _slide_text(slide)


def test_png_size_dpi_and_display_height(tmp_path):
    import pymupdf

    from app.pptx_export import CONTENT_W, _image_display_height, _png_size_dpi

    pix = pymupdf.Pixmap(pymupdf.csRGB, pymupdf.IRect(0, 0, 1920, 480))
    pix.clear_with(200)
    pix.set_dpi(96, 96)
    p = tmp_path / "wide.png"
    pix.save(str(p))

    w, h, dpi = _png_size_dpi(str(p))
    assert (w, h) == (1920, 480)
    assert 94 <= dpi <= 98  # ~96, allowing rounding via the metres conversion

    # 1920px / 96dpi = 20in wide > the column, so it scales down keeping aspect.
    native_w = int(1920 / dpi * 914400)
    native_h = int(480 / dpi * 914400)
    expected = int(native_h * CONTENT_W / native_w)
    assert abs(_image_display_height(str(p), CONTENT_W) - expected) < 6000

    assert _png_size_dpi(str(tmp_path / "missing.png")) is None
    not_png = tmp_path / "fake.png"
    not_png.write_bytes(b"not a png at all")
    assert _png_size_dpi(str(not_png)) is None


def test_export_long_section_paginates(tmp_path):
    body = "\n\n".join(f"Paragraph number {i} with some words." for i in range(60))
    md = f"## Long\n\n{body}\n"
    out = tmp_path / "long.pptx"
    count = export_markdown_to_pptx(md, out)
    # One H2 section, but it overflows a single slide -> continuation slides.
    assert count >= 2
    assert len(Presentation(str(out)).slides) == count
