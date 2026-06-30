"""Export a Markdown document to an editable PowerPoint (.pptx) deck.

No Office / LibreOffice needed — ``python-pptx`` writes the OOXML directly, the
same spirit as PDFium for the PDF side. The Markdown is split into slides (by
``---`` thematic breaks if the document has any, otherwise by the ``##`` heading
level) and mapped to *native* PowerPoint objects — title, bullet lists, code
boxes, tables and images — so the result stays editable in PowerPoint.

Mermaid diagrams and ``$$`` math are rendered to images when an *image_provider*
is supplied (the GUI wires one up using the web engine); without one — or if a
fragment fails to render — they degrade to a labelled source box.

The parsing / slide-splitting logic is pure (no Qt), so it is unit-testable by
generating a deck and reading it back with python-pptx.
"""

from __future__ import annotations

import io
import struct
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path

from markdown_it import MarkdownIt

try:  # math support matches the viewer's renderer (already a dependency)
    from mdit_py_plugins.dollarmath import dollarmath_plugin
except Exception:  # pragma: no cover - optional
    dollarmath_plugin = None

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ----------------------------- document model -----------------------------


@dataclass
class Run:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


@dataclass
class Heading:
    level: int
    runs: list


@dataclass
class Para:
    runs: list


@dataclass
class ListItem:
    level: int
    runs: list
    ordered: bool = False
    number: int = 1


@dataclass
class ListBlock:
    items: list


@dataclass
class Code:
    text: str
    lang: str = ""


@dataclass
class Table:
    header: list  # list[str]
    rows: list  # list[list[str]]


@dataclass
class Image:
    src: str
    alt: str = ""


@dataclass
class Quote:
    runs: list


class _Hr:
    pass


HR = _Hr()


@dataclass
class SlideModel:
    title: list | None = None  # list[Run] or None
    blocks: list = field(default_factory=list)


# ------------------------------- parsing ---------------------------------


def _build_parser() -> MarkdownIt:
    md = MarkdownIt("commonmark", {"html": False, "linkify": False})
    md.enable("table")
    try:
        md.enable("strikethrough")
    except Exception:  # pragma: no cover - depends on preset
        pass
    if dollarmath_plugin:
        md.use(dollarmath_plugin)
    return md


def _runs(inline) -> list:
    """Flatten an inline token's children into styled text runs."""
    if inline is None:
        return []
    children = inline.children
    if not children:
        return [Run(inline.content)] if inline.content else []
    runs: list = []
    bold = italic = 0
    for c in children:
        t = c.type
        if t == "text":
            runs.append(Run(c.content, bold > 0, italic > 0))
        elif t == "code_inline":
            runs.append(Run(c.content, bold > 0, italic > 0, code=True))
        elif t == "math_inline":
            runs.append(Run(c.content, bold > 0, italic > 0, code=True))
        elif t == "strong_open":
            bold += 1
        elif t == "strong_close":
            bold = max(0, bold - 1)
        elif t == "em_open":
            italic += 1
        elif t == "em_close":
            italic = max(0, italic - 1)
        elif t in ("softbreak", "hardbreak"):
            runs.append(Run(" "))
        # link_open/close: keep the visible text (children carry it); image
        # tokens are pulled out separately by _images().
    return [r for r in runs if r.text]


def _plain(inline) -> str:
    return "".join(r.text for r in _runs(inline))


def _images(inline) -> list:
    if inline is None or not inline.children:
        return []
    out = []
    for c in inline.children:
        if c.type == "image":
            src = c.attrs.get("src", "") if hasattr(c, "attrs") else ""
            out.append(Image(src, c.content or ""))
    return out


def _parse_list(tokens, i, level):
    open_tok = tokens[i]
    ordered = open_tok.type == "ordered_list_open"
    items: list = []
    number = 1
    i += 1
    while i < len(tokens) and tokens[i].type not in (
        "bullet_list_close",
        "ordered_list_close",
    ):
        if tokens[i].type == "list_item_open":
            i += 1
            item_runs: list = []
            nested: list = []
            while i < len(tokens) and tokens[i].type != "list_item_close":
                tt = tokens[i].type
                if tt == "paragraph_open":
                    if item_runs:
                        item_runs.append(Run(" "))
                    item_runs.extend(_runs(tokens[i + 1]))
                    i += 3
                elif tt == "inline":
                    item_runs.extend(_runs(tokens[i]))
                    i += 1
                elif tt in ("bullet_list_open", "ordered_list_open"):
                    sub, i = _parse_list(tokens, i, level + 1)
                    nested.extend(sub.items)
                else:
                    i += 1
            items.append(ListItem(level, item_runs, ordered, number))
            number += 1
            items.extend(nested)
            i += 1  # consume list_item_close
        else:
            i += 1
    i += 1  # consume the list close token
    return ListBlock(items), i


def _parse_table(tokens, i):
    header: list = []
    rows: list = []
    in_head = False
    current: list | None = None
    i += 1  # consume table_open
    while i < len(tokens) and tokens[i].type != "table_close":
        tt = tokens[i].type
        if tt == "thead_open":
            in_head = True
            i += 1
        elif tt == "thead_close":
            in_head = False
            i += 1
        elif tt in ("tbody_open", "tbody_close"):
            i += 1
        elif tt == "tr_open":
            current = []
            i += 1
        elif tt == "tr_close":
            if in_head:
                header = current or []
            else:
                rows.append(current or [])
            current = None
            i += 1
        elif tt in ("th_open", "td_open"):
            if current is not None:
                current.append(_plain(tokens[i + 1]))
            i += 3  # open, inline, close
        else:
            i += 1
    i += 1  # consume table_close
    return Table(header, rows), i


def _parse_blockquote(tokens, i):
    depth = 1
    runs: list = []
    i += 1  # consume blockquote_open
    while i < len(tokens) and depth > 0:
        tt = tokens[i].type
        if tt == "blockquote_open":
            depth += 1
            i += 1
        elif tt == "blockquote_close":
            depth -= 1
            i += 1
        elif tt == "paragraph_open":
            if runs:
                runs.append(Run(" "))
            runs.extend(_runs(tokens[i + 1]))
            i += 3
        elif tt == "inline":
            runs.extend(_runs(tokens[i]))
            i += 1
        else:
            i += 1
    return Quote(runs), i


def parse_elements(md_text: str) -> list:
    """Parse Markdown into a flat list of block elements."""
    md = _build_parser()
    tokens = md.parse(md_text)
    els: list = []
    i = 0
    n = len(tokens)
    while i < n:
        t = tokens[i]
        tt = t.type
        if tt == "heading_open":
            level = int(t.tag[1])
            els.append(Heading(level, _runs(tokens[i + 1])))
            i += 3
        elif tt == "paragraph_open":
            inline = tokens[i + 1]
            imgs = _images(inline)
            text_runs = _runs(inline)
            if imgs:
                els.extend(imgs)
                if text_runs:
                    els.append(Para(text_runs))
            elif text_runs:
                els.append(Para(text_runs))
            i += 3
        elif tt in ("fence", "code_block"):
            els.append(Code(t.content.rstrip("\n"), (t.info or "").strip()))
            i += 1
        elif tt == "math_block":
            els.append(Code(t.content.strip(), "math"))
            i += 1
        elif tt == "hr":
            els.append(HR)
            i += 1
        elif tt in ("bullet_list_open", "ordered_list_open"):
            block, i = _parse_list(tokens, i, 0)
            els.append(block)
        elif tt == "table_open":
            block, i = _parse_table(tokens, i)
            els.append(block)
        elif tt == "blockquote_open":
            block, i = _parse_blockquote(tokens, i)
            els.append(block)
        else:
            i += 1
    return els


# --------------------------- slide splitting -----------------------------


def _first_heading_title(blocks):
    for idx, b in enumerate(blocks):
        if isinstance(b, Heading):
            return b.runs, blocks[:idx] + blocks[idx + 1 :]
    return None, blocks


def split_into_slides(elements: list) -> list:
    """Split elements into slides: by ``---`` if present, else by heading."""
    if any(e is HR for e in elements):
        slides = []
        segment: list = []
        for e in elements:
            if e is HR:
                title, blocks = _first_heading_title(segment)
                slides.append(SlideModel(title, blocks))
                segment = []
            else:
                segment.append(e)
        title, blocks = _first_heading_title(segment)
        slides.append(SlideModel(title, blocks))
        return [s for s in slides if s.title or s.blocks]

    levels = [e.level for e in elements if isinstance(e, Heading)]
    if not levels:
        return [SlideModel(None, list(elements))] if elements else []
    split_level = 2 if 2 in levels else min(levels)

    slides = []
    current = SlideModel(None, [])
    started = False
    for e in elements:
        if isinstance(e, Heading) and e.level <= split_level:
            if started and (current.title or current.blocks):
                slides.append(current)
            current = SlideModel(e.runs, [])
            started = True
        else:
            current.blocks.append(e)
    if current.title or current.blocks:
        slides.append(current)
    return slides


# ------------------------------ rendering --------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_TOP = Inches(0.5)
BLOCK_GAP = Inches(0.14)
TITLE_H = Inches(0.95)

TITLE_PT = 30
HEAD_PT = 20
BODY_PT = 16
CODE_PT = 12
TABLE_PT = 12

_CODE_BG = RGBColor(0xF2, 0xF2, 0xF2)
_CODE_FG = RGBColor(0x1A, 0x1A, 0x1A)
_QUOTE_FG = RGBColor(0x55, 0x55, 0x55)
_TITLE_FG = RGBColor(0x1F, 0x1F, 0x1F)


def _emit_runs(paragraph, runs, size_pt, color=None):
    if not runs:
        runs = [Run("")]
    for r in runs:
        run = paragraph.add_run()
        run.text = r.text
        run.font.size = Pt(size_pt)
        if r.bold:
            run.font.bold = True
        if r.italic:
            run.font.italic = True
        if r.code:
            run.font.name = "Consolas"
        if color is not None:
            run.font.color.rgb = color


def _disp_width(text: str) -> int:
    return sum(2 if ord(c) > 0x2E7F else 1 for c in text)


def _est_lines(text: str, font_pt: int, width_emu: int) -> int:
    width_in = width_emu / 914400.0
    char_w_in = 0.55 * font_pt / 72.0
    cpl = max(1, int(width_in / char_w_in))
    total = 0
    for line in (text or "").split("\n"):
        total += max(1, -(-_disp_width(line) // cpl))
    return max(1, total)


def _line_h(font_pt: int) -> int:
    return Emu(int(font_pt * 1.35 / 72.0 * 914400))


def _est_block_height(block) -> int:
    if isinstance(block, Heading):
        return _est_lines(_text_of(block.runs), HEAD_PT, CONTENT_W) * _line_h(HEAD_PT) + Inches(0.1)
    if isinstance(block, Para):
        return _est_lines(_text_of(block.runs), BODY_PT, CONTENT_W) * _line_h(BODY_PT) + Inches(0.1)
    if isinstance(block, Quote):
        return _est_lines(_text_of(block.runs), BODY_PT, CONTENT_W) * _line_h(BODY_PT) + Inches(0.2)
    if isinstance(block, ListBlock):
        h = Inches(0.1)
        for it in block.items:
            h += _est_lines(_text_of(it.runs), BODY_PT, CONTENT_W) * _line_h(BODY_PT)
        return h
    if isinstance(block, Code):
        lines = len(block.text.split("\n"))
        line_est = Emu(lines * _line_h(CODE_PT)) + Inches(0.25)
        # A mermaid/math block may become a rendered image of unknown height;
        # reserve a generous default so pagination leaves room for it.
        if block.lang == "mermaid":
            return max(line_est, Inches(3.0))
        if block.lang == "math":
            return max(line_est, Inches(0.7))
        return line_est
    if isinstance(block, Table):
        return (len(block.rows) + 1) * Inches(0.4)
    if isinstance(block, Image):
        return Inches(3.2)  # rough; actual measured at render time
    return Inches(0.4)


def _text_of(runs) -> str:
    return "".join(r.text for r in runs)


def _bullet_prefix(item: ListItem) -> str:
    indent = "    " * item.level
    if item.ordered:
        return f"{indent}{item.number}. "
    glyph = ["•", "◦", "‣", "·"][min(item.level, 3)]
    return f"{indent}{glyph} "


def _render_block(
    slide, block, left, top, width, base_dir, image_provider=None, prerendered=None
) -> int:
    """Render one block at (left, top); return the height it consumed (EMU)."""
    if isinstance(block, (Heading, Para, Quote)):
        is_quote = isinstance(block, Quote)
        size = HEAD_PT if isinstance(block, Heading) else BODY_PT
        height = _est_block_height(block)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        color = _QUOTE_FG if is_quote else None
        _emit_runs(p, block.runs, size, color)
        if isinstance(block, Heading):
            for run in p.runs:
                run.font.bold = True
        if is_quote:
            for run in p.runs:
                run.font.italic = True
        return height

    if isinstance(block, ListBlock):
        height = _est_block_height(block)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(block.items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            prefix = Run(_bullet_prefix(item))
            _emit_runs(p, [prefix] + item.runs, BODY_PT)
        return height

    if isinstance(block, Code):
        # Mermaid diagrams and LaTeX math render to images. _render_slide does
        # the provider call (so it can size pagination from the real image) and
        # passes the PNG path here; on failure fall through to the source box.
        if block.lang in ("mermaid", "math") and image_provider is not None:
            if prerendered:
                return _render_image(
                    slide, Image(prerendered, block.lang), left, top, width, None
                )
        lines = block.text.split("\n")
        height = Emu(len(lines) * _line_h(CODE_PT)) + Inches(0.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _CODE_BG
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        if block.lang in ("mermaid", "math"):
            label = "Mermaid 圖（原始碼）" if block.lang == "mermaid" else "數學式（LaTeX）"
            head = tf.paragraphs[0]
            r = head.add_run()
            r.text = label
            r.font.size = Pt(CODE_PT - 1)
            r.font.italic = True
            r.font.color.rgb = _QUOTE_FG
            first = tf.add_paragraph()
        else:
            first = tf.paragraphs[0]
        for idx, line in enumerate(lines):
            p = first if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.name = "Consolas"
            run.font.size = Pt(CODE_PT)
            run.font.color.rgb = _CODE_FG
        return height

    if isinstance(block, Table):
        n_cols = max(
            [len(block.header)] + [len(r) for r in block.rows] or [1]
        ) or 1
        n_rows = len(block.rows) + (1 if block.header else 0)
        n_rows = max(1, n_rows)
        row_h = Inches(0.4)
        height = n_rows * row_h
        gtable = slide.shapes.add_table(
            n_rows, n_cols, left, top, width, height
        ).table
        data_rows = ([block.header] if block.header else []) + block.rows
        for r_idx, row in enumerate(data_rows):
            for c_idx in range(n_cols):
                cell = gtable.cell(r_idx, c_idx)
                cell.text = row[c_idx] if c_idx < len(row) else ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(TABLE_PT)
                        if block.header and r_idx == 0:
                            run.font.bold = True
        return height

    if isinstance(block, Image):
        return _render_image(slide, block, left, top, width, base_dir)

    return 0


def _load_image_stream(src: str, base_dir):
    """Return a BytesIO for *src* (local path or remote URL), or None."""
    if src.startswith(("http://", "https://")):
        try:
            with urllib.request.urlopen(src, timeout=8) as resp:  # nosec - user content
                return io.BytesIO(resp.read())
        except Exception:
            return None
    path = Path(src)
    if not path.is_absolute() and base_dir is not None:
        path = Path(base_dir) / src
    try:
        return io.BytesIO(path.read_bytes())
    except Exception:
        return None


def _render_image(slide, block, left, top, width, base_dir) -> int:
    stream = _load_image_stream(block.src, base_dir)
    if stream is None:
        # Fall back to a labelled placeholder so the slide still reads sensibly.
        box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"[圖片：{block.alt or block.src}]"
        r.font.size = Pt(BODY_PT)
        r.font.italic = True
        r.font.color.rgb = _QUOTE_FG
        return Inches(0.4)
    pic = slide.shapes.add_picture(stream, left, top)
    w, h = pic.width, pic.height
    if w > width:
        h = int(h * width / w)
        w = width
    max_h = SLIDE_H - MARGIN - top
    if h > max_h and h > 0:
        scale = max_h / h
        w = int(w * scale)
        h = int(h * scale)
    pic.width = w
    pic.height = h
    return h


def _add_title(slide, runs, continuation: bool) -> int:
    box = slide.shapes.add_textbox(MARGIN, CONTENT_TOP, CONTENT_W, TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    title_runs = list(runs)
    if continuation:
        title_runs = title_runs + [Run("（續）")]
    _emit_runs(p, title_runs, TITLE_PT, _TITLE_FG)
    for run in p.runs:
        run.font.bold = True
    return TITLE_H


def _png_size_dpi(path):
    """Return (width_px, height_px, dpi) for a PNG, or None. Pure stdlib."""
    try:
        with open(path, "rb") as f:
            blob = f.read(256)
    except OSError:
        return None
    if blob[:8] != b"\x89PNG\r\n\x1a\n" or len(blob) < 24:
        return None
    w, h = struct.unpack(">II", blob[16:24])
    dpi = 72
    idx = blob.find(b"pHYs")
    if idx != -1 and idx + 13 <= len(blob):
        ppu_x, _ppu_y, unit = struct.unpack(">IIB", blob[idx + 4 : idx + 13])
        if unit == 1 and ppu_x:  # pixels per metre
            dpi = max(1, round(ppu_x * 0.0254))
    return w, h, dpi


def _image_display_height(path, width_emu) -> int:
    """Height (EMU) a PNG will occupy after scaling to fit *width_emu* wide."""
    info = _png_size_dpi(path)
    if not info:
        return Inches(3.0)
    pw, ph, dpi = info
    w_emu = int(pw / dpi * 914400)
    h_emu = int(ph / dpi * 914400)
    if w_emu > width_emu and w_emu > 0:
        h_emu = int(h_emu * width_emu / w_emu)
    return h_emu


def _render_slide(prs, layout, model: SlideModel, base_dir, image_provider=None) -> int:
    """Render one slide model into the deck, paginating long content. Returns
    the number of PowerPoint slides actually produced."""
    bottom = SLIDE_H - MARGIN
    produced = 0

    def new_slide(continuation):
        nonlocal produced
        slide = prs.slides.add_slide(layout)
        produced += 1
        top = CONTENT_TOP
        if model.title:
            top += _add_title(slide, model.title, continuation) + BLOCK_GAP
        return slide, top, top

    slide, cursor, first_top = new_slide(False)
    for block in model.blocks:
        # Render mermaid/math up front (cached, so it's not re-rendered later)
        # so pagination can use the image's real height instead of a guess.
        prerendered = None
        if (
            image_provider is not None
            and isinstance(block, Code)
            and block.lang in ("mermaid", "math")
        ):
            prerendered = image_provider(block.lang, block.text)
        if prerendered:
            est = _image_display_height(prerendered, CONTENT_W) + Inches(0.1)
        else:
            est = _est_block_height(block)
        if cursor > first_top and cursor + est > bottom:
            slide, cursor, first_top = new_slide(True)
        used = _render_block(
            slide, block, MARGIN, cursor, CONTENT_W, base_dir,
            image_provider, prerendered,
        )
        cursor += used + BLOCK_GAP
    return produced


def export_markdown_to_pptx(md_text: str, out_path, base_dir=None, image_provider=None) -> int:
    """Convert *md_text* to a .pptx deck at *out_path*. Returns slide count.

    *image_provider*, when given, is ``(kind, source) -> png_path | None`` for
    ``kind`` in {"mermaid", "math"}; it renders those fragments to images that
    are embedded instead of source-code boxes. ``None`` (or a ``None`` return)
    falls back to the labelled source box, keeping this module Qt-free.
    """
    elements = parse_elements(md_text)
    slides = split_into_slides(elements)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # Blank layout — we place every shape ourselves.

    total = 0
    if not slides:
        prs.slides.add_slide(blank)
        total = 1
    for model in slides:
        total += _render_slide(prs, blank, model, base_dir, image_provider)

    prs.save(str(out_path))
    return total
